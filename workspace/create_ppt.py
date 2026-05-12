#!/usr/bin/env python
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

PRIMARY_COLOR = RGBColor(31, 73, 125)
ACCENT_COLOR = RGBColor(79, 129, 189)
TEXT_DARK = RGBColor(0, 0, 0)


def create_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)
    return slide


def create_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True

    tf = slide.placeholders[1].text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(12)

    return slide


def main():
    output_path = "D:/pythonProject/or_llm_agent/workspace/中期答辩_黄纤纤.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    create_title_slide(
        prs,
        "基于大模型的生鲜物流配送路径优化研究",
        "本科毕业设计（论文）中期答辩\n\n学生姓名：黄纤纤\n学号：2022216431\n指导教师：程八一 教授\n\n2026年5月",
    )

    create_content_slide(
        prs,
        "汇报提纲",
        [
            "一、研究背景与问题概述",
            "二、第三章：生鲜物流配送路径优化模型构建",
            "  - 问题描述",
            "  - 参数符号说明",
            "  - 目标函数与约束条件",
            "三、第四章：基于大模型与ALNS算法的问题求解",
            "  - ALNS求解框架设计",
            "  - 基于大模型的启发式算子生成机制",
            "  - Prompt工程设计",
            "四、实验设计与预期成果",
            "五、总结与后续工作计划",
        ],
    )

    create_content_slide(
        prs,
        "一、研究背景与问题概述",
        [
            "生鲜物流配送路径问题（VRPTW-FC）是经典VRP的扩展变体",
            "需同时考虑：",
            "  - 车辆容量约束",
            "  - 时间窗约束（软/硬）",
            "  - 冷链制冷成本",
            "  - 生鲜产品新鲜度衰减",
            "传统ALNS算法依赖人工经验设计算子，耗时且难以自适应",
        ],
    )

    create_content_slide(
        prs,
        "二、第三章：生鲜物流配送路径优化模型构建",
        [
            "3.1 问题描述",
            "    带时间窗约束的生鲜物流车辆路径问题（VRPTW-FC）",
            "3.2 模型构建",
            "    - 参数符号说明",
            "    - 目标函数（五项成本）",
            "    - 约束条件",
        ],
    )

    create_content_slide(
        prs,
        "3.1 问题描述",
        [
            "假设某生鲜配送中心配备m辆同质配送车辆，需向n个客户配送生鲜产品",
            "每个客户具备：",
            "  - 需求量",
            "  - 地理坐标",
            "  - 服务时间���（软时间窗 + 硬时间窗）",
            "额外考虑因素：",
            "  - 冷链制冷成本",
            "  - 货物新鲜度衰减（指数型）",
            "  - 软硬双层时间窗约束",
        ],
    )

    create_content_slide(
        prs,
        "3.2.1 参数符号说明",
        [
            "集合：V-节点集（含配送中心0），K-车辆集",
            "决策变量：x_ijk-车辆k从i到j，y_jk-客户j被车辆k服务",
            "关键参数：",
            "  - Q_k：车辆最大装载量",
            "  - d_ij：客户i与j之间的欧氏距离",
            "  - c_cold：单位时间制冷成本",
            "  - alpha, beta：新鲜度衰减系数",
            "  - TW_j^soft：客户期望时间窗（软）",
            "  - TW_j^hard：客户可接受时间窗（硬）",
        ],
    )

    create_content_slide(
        prs,
        "3.2.2 目标函数",
        [
            "综合配送总成本 = 车辆固定成本 + 运输距离成本",
            "                + 冷链制冷成本 + 货损成本 + 时间窗惩罚成本",
            " ",
            "(1) 车辆固定成本 C_f = m * c_fixed",
            "(2) 运输距离成本 C_d = sum(d_ij * c_dist)",
            "(3) 冷链制冷成本 C_cold = T_total * c_cold",
            "(4) 货损成本 C_damage = sum(q_j * p * decay_j)",
            "(5) 时间窗惩罚成本 C_TW = 分段线性函数",
        ],
    )

    create_content_slide(
        prs,
        "3.2.3 约束条件",
        [
            "约束(1) 客户访问约束：每个客户恰好被一辆车服务一次",
            "约束(2) 流量守恒约束：车辆驶入客户节点后必须驶出",
            "约束(3) 起讫点约束：车辆从配送中心出发并最终返回",
            "约束(4) 载量约束：车辆载重量不超过最大装载量",
            "约束(5) 时间窗约束：服务时间必须在硬时间窗内",
            "约束(6) 可用车辆数约束：实际使用车辆数不超过可用数",
        ],
    )

    create_content_slide(
        prs,
        "三、第四章：基于大模型与ALNS算法的问题求解",
        [
            "4.1 基于大模型的启发式算子生成机制",
            "    - 算子自动化生成框架设计",
            "    - Prompt工程设计",
            "4.2 ALNS求解框架设计",
            "    - 初始解构造策略",
            "    - 破坏算子设计",
            "    - 修复算子设计",
            "    - 自适应权重与破坏比例更新机制",
        ],
    )

    create_content_slide(
        prs,
        "4.2 ALNS求解框架设计",
        [
            "采用骨架-插件分离架构",
            "五大核心模块：",
            "  (1) 数据预处理与初始化模块",
            "  (2) 算子注册与调度模块（轮盘赌选择）",
            "  (3) 容错与降级模块（备用算子）",
            "  (4) 解评估与验证模块",
            "  (5) 迭代控制模块",
            " ",
            "破坏算子：随机移除、路径移除、相关性移除",
            "修复算子：贪心插入、后悔插入、随机插入",
        ],
    )

    create_content_slide(
        prs,
        "4.1 基于大模型的启发式算子优化机制",
        [
            "核心思想：基线评估-反馈生成-动态注入-贪心接受",
            " ",
            "实现流程：",
            " 1. 运行基线ALNS获取初始解及多维性能指标",
            " 2. 将性能数据结构化嵌入Prompt",
            " 3. LLM分析瓶颈并生成改进算子",
            " 4. 动态代码注入嵌入求解器",
            " 5. 验证并决定接受或回退",
            " ",
            "性能指标：全局成本、收敛速度、算子成功率、成本结构",
        ],
    )

    create_content_slide(
        prs,
        "4.1.2 Prompt工程设计",
        [
            "结构化Prompt构建方法（五���次）：",
            " 1. 角色设定：明确算子优化任务目标",
            " 2. 问题建模：描述问题特征与约束",
            " 3. 系统约束：定义输出格式与方法签名",
            " 4. 状态输入：嵌入性能指标数据",
            " 5. 优化引导：引导模型进行策略分析",
            " ",
            "引入历史优化记录机制：",
            " - 避免重复策略",
            " - 增强策略多样性",
        ],
    )

    create_content_slide(
        prs,
        "四、实验设计",
        [
            "数据集：Solomon VRP Benchmark Dataset",
            " - 25客户规模（聚类型/随机型/混合型）",
            " - 50客户规模",
            " - 75客户规模",
            " ",
            "算法参数配置：",
            " - 迭代次数：500-1000次",
            " - 破坏比例：5%-30%",
            " - 算子集合：5个破坏+5个修复",
            " - 优化轮次：5轮",
            " ",
            "对比算法：传统ALNS、GA、SA",
        ],
    )

    create_content_slide(
        prs,
        "五、预期成果",
        [
            "创新点：",
            " 1. 提出基于LLM的ALNS算子自动化优化框架",
            " 2. 设计多维度性能评价指标体系",
            " 3. 构建结构化Prompt工程方法",
            " 4. 实现算子代码动态注入机制",
            " ",
            "预期效果：",
            " - 收敛速度提升15%以上",
            " - 解质量改善10%以上",
            " - 算子设计时间减少50%以上",
        ],
    )

    create_content_slide(
        prs,
        "六、总结与后续计划",
        [
            "已完成：",
            " - 第三章：生鲜物流配送路径优化模型构建",
            " - 第四章：基于大模型与ALNS算法框架设计",
            " ",
            "后续工作：",
            " - 完成实验验证与结果分析",
            " - 完善论文撰写",
            " - 准备毕业答辩",
        ],
    )

    create_title_slide(prs, "谢谢！", "恳请各位老师批评指正")

    prs.save(output_path)
    print("PPT saved to: " + output_path)


if __name__ == "__main__":
    main()
