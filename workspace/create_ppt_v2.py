#!/usr/bin/env python
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import nsmap


def create_title_slide_v2(prs, title, subtitle, date_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Blue rectangle for title area
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.8), Inches(13.333), Inches(1.5)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(4, 64, 126)
    title_bar.line.fill.background()

    # Title text on blue bar
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(2), Inches(12.7), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = 1

    # Subtitle (decorative text)
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(3.8), Inches(12.7), Inches(0.6)
    )
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "本科毕业设计（论文）中期答辩"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(31, 73, 125)
    p2.alignment = 1

    # Info at bottom left
    info_box1 = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.3), Inches(4), Inches(0.4)
    )
    tf3 = info_box1.text_frame
    tf3.paragraphs[0].text = "学生姓名：黄纤纤"
    tf3.paragraphs[0].font.size = Pt(18)
    tf3.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    info_box2 = slide.shapes.add_textbox(
        Inches(4.5), Inches(6.3), Inches(4), Inches(0.4)
    )
    tf4 = info_box2.text_frame
    tf4.paragraphs[0].text = "学号：2022216431"
    tf4.paragraphs[0].font.size = Pt(18)
    tf4.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    info_box3 = slide.shapes.add_textbox(
        Inches(8.5), Inches(6.3), Inches(4), Inches(0.4)
    )
    tf5 = info_box3.text_frame
    tf5.paragraphs[0].text = "指导教师：程八一 教授"
    tf5.paragraphs[0].font.size = Pt(18)
    tf5.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Date at bottom right
    date_box = slide.shapes.add_textbox(Inches(10), Inches(6.3), Inches(3), Inches(0.4))
    tf6 = date_box.text_frame
    tf6.paragraphs[0].text = date_text
    tf6.paragraphs[0].font.size = Pt(18)
    tf6.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    tf6.paragraphs[0].alignment = 2

    return slide


def create_content_slide_v2(prs, title, bullets=None, section_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Style title - blue background bar
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Get the title shape background
    for shape in slide.shapes:
        if shape == title_shape:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(4, 64, 126)
            shape.line.fill.background()
            break

    # Content bullets
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    if bullets:
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.space_before = Pt(10)

            # Indent sub-bullets
            if bullet.startswith("  ") or bullet.startswith("- "):
                p.level = 1
                p.font.size = Pt(18)

    return slide


def create_section_slide(prs, section_title, section_num):
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Blue header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(4, 64, 126)
    header.line.fill.background()

    # Section number
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(1), Inches(0.6))
    tf = num_box.text_frame
    tf.paragraphs[0].text = section_num
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(0.3), Inches(10), Inches(0.6)
    )
    tf2 = title_box.text_frame
    tf2.paragraphs[0].text = section_title
    tf2.paragraphs[0].font.size = Pt(32)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    return slide


def main():
    output_path = "D:/pythonProject/or_llm_agent/workspace/中期答辩_黄纤纤_v2.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    create_title_slide_v2(
        prs,
        "基于大模型的生鲜物流配送路径优化研究",
        "本科毕业设计（论文）中期答辩",
        "2026年5月",
    )

    # Slide 2: Contents
    create_content_slide_v2(
        prs,
        "汇报提纲",
        [
            "一、研究背景与问题概述",
            "二、第三章：生鲜物流配送路径优化模型构建",
            "  - 问题描述与参数符号",
            "  - 目标函数与约束条件",
            "三、第四章：基于大模型与ALNS算法的问题求解",
            "  - ALNS求解框架设计",
            "  - 基于大模型的启发式算子优化机制",
            "  - Prompt工程设计",
            "四、实验设计与预期成果",
            "五、总结与后续计划",
        ],
    )

    # Slide 3: Background
    create_content_slide_v2(
        prs,
        "一、研究背景与问题概述",
        [
            "生鲜物流配送路径问题（VRPTW-FC）是经典VRP的扩展变体",
            "需要同时考虑多种约束：",
            "- 车辆容量限制",
            "- 时间窗约束（软/硬双层）",
            "- 冷链制冷成本",
            "- 生鲜产品新鲜度衰减（指数型）",
            "传统ALNS算法依赖人工经验设计算子，效率低且难以自适应",
        ],
    )

    # Slide 4: Chapter 3 Overview
    create_content_slide_v2(
        prs,
        "二、第三章：生鲜物流配送路径优化模型构建",
        [
            "3.1 问题描述",
            "  带时间窗约束的生鲜物流车辆路径问题（VRPTW-FC）",
            "3.2 模型构建",
            "  - 参数符号说明",
            "  - 目标函数（五项成本）",
            "  - 约束条件",
        ],
    )

    # Slide 5: Problem Description
    create_content_slide_v2(
        prs,
        "3.1 问题描述",
        [
            "场景假设：生鲜配送中心配备m辆同质配送车辆，需向n个客户配送生鲜产品",
            "客户属性：",
            "- 需求量 q_j",
            "- 地理坐标 (x_j, y_j)",
            "- 服务时间窗 TW_j（软时间窗 + 硬时间窗）",
            "额外考虑因素：",
            "- 冷链制冷成本：与运输时间成正比",
            "- 货物新鲜度衰减：指数型衰减模型",
            "- 软硬双层时间窗：软窗外产生惩罚，硬窗外违约",
        ],
    )

    # Slide 6: Parameters
    create_content_slide_v2(
        prs,
        "3.2.1 参数符号说明",
        [
            "集合定义：",
            "- V：节点集（含配送中心0）",
            "- K：车辆集",
            "决策变量：",
            "- x_ijk：车辆k从i行驶到j（0-1变量）",
            "- y_jk：客户j被车辆k服务（0-1变量）",
            "关键参数：",
            "- Q_k：车辆最大装载量",
            "- d_ij：客户i与j之间的欧氏距离",
            "- c_cold：单位时间制冷成本",
            "- alpha, beta：新鲜度衰减系数",
        ],
    )

    # Slide 7: Objective Function
    create_content_slide_v2(
        prs,
        "3.2.2 目标函数",
        [
            "综合配送总成本 = 车辆固定成本 + 运输距离成本",
            "                + 冷链制冷成本 + 货损成本 + 时间窗惩罚成本",
            " ",
            "(1) 车辆固定成本 C_f = m * c_fixed",
            "(2) 运输距离成本 C_d = sum(d_ij * c_dist)",
            "(3) 冷链制冷成本 C_cold = T_total * c_cold",
            "(4) 货损成本 C_damage = sum(q_j * p * decay_j)",
            "(5) 时间窗惩罚成本 C_TW = 分段线性惩罚函数",
        ],
    )

    # Slide 8: Constraints
    create_content_slide_v2(
        prs,
        "3.2.3 约束条件",
        [
            "约束(1) 客户访问约束：每个客户恰好被一辆车服务一次",
            "约束(2) 流量守恒约束：车辆驶入客户节点后必须驶出",
            "约束(3) 起讫点约束：从配送中心出发并最终返回",
            "约束(4) 容量约束：车辆载重量不超过最大装载量",
            "约束(5) 时间窗约束：服务时间必须在硬时间窗内",
            "约束(6) 车辆数约束：实际使用车辆数不超过可用数",
        ],
    )

    # Slide 9: Chapter 4 Overview
    create_content_slide_v2(
        prs,
        "三、第四章：基于大模型与ALNS算法的问题求解",
        [
            "4.1 基于大模型的启发式算子生成机制",
            "  - 算子自动化生成框架设计",
            "  - Prompt工程设计",
            "4.2 ALNS求解框架设计",
            "  - 初始解构造策略",
            "  - 破坏算子设计",
            "  - 修复算子设计",
            "  - 自适应权重更新机制",
        ],
    )

    # Slide 10: ALNS Framework
    create_content_slide_v2(
        prs,
        "4.2 ALNS求解框架设计",
        [
            "采用骨架-插件分离架构",
            "五大核心模块：",
            "(1) 数据预处理与初始化模块",
            "(2) 算子注册与调度模块（轮盘赌选择）",
            "(3) 容错与降级模块（备用算子）",
            "(4) 解评估与验证模块",
            "(5) 迭代控制模块",
            " ",
            "破坏算子：随机移除、路径移除、相关性移除",
            "修复算子：贪心插入、后悔插入、随机插入",
        ],
    )

    # Slide 11: LLM Operator Optimization
    create_content_slide_v2(
        prs,
        "4.1 基于大模型的启发式算子优化机制",
        [
            "核心思想：基线评估-反馈生成-动态注入-贪心接受",
            " ",
            "实现流程：",
            "1. 运行基线ALNS获取初始解及多维性能指标",
            "2. 将性能数据结构化嵌入Prompt",
            "3. LLM分析瓶颈并生成改进算子",
            "4. 动态代码注入嵌入求解器",
            "5. 验证并决定接受或回退",
            " ",
            "性能指标：全局成本、收敛速度、算子成功率、��本��构",
        ],
    )

    # Slide 12: Prompt Engineering
    create_content_slide_v2(
        prs,
        "4.1.2 Prompt工程设计",
        [
            "结构化Prompt构建方法（五层次）：",
            "1. 角色设定：明确算子优化任务目标",
            "2. 问题建模：描述问题特征与约束",
            "3. 系统约束：定义输出格式与方法签名",
            "4. 状态输入：嵌入性能指标数据",
            "5. 优化引导：引导模型进行策略分析",
            " ",
            "历史优化记录机制：",
            "- 避免重复策略",
            "- 增强策略多样性",
        ],
    )

    # Slide 13: Experimental Design
    create_content_slide_v2(
        prs,
        "四、实验设计",
        [
            "数据集：Solomon VRP Benchmark Dataset",
            "- 25客户规模（聚类型/随机型/混合型）",
            "- 50客户规模",
            "- 75客户规模",
            " ",
            "算法参数配置：",
            "- 迭代次数：500-1000次",
            "- 破坏比例：5%-30%",
            "- 算子集合：5个破坏+5个修复",
            "- 优化轮次：5轮",
            " ",
            "对比算法：传统ALNS、GA、SA",
        ],
    )

    # Slide 14: Expected Results
    create_content_slide_v2(
        prs,
        "五、预期成果",
        [
            "创新点：",
            "1. 提出基于LLM的ALNS算子自动化优化框架",
            "2. 设计多维度性能评价指标体系",
            "3. 构建结构化Prompt工程方法",
            "4. 实现算子代码动态注入机制",
            " ",
            "预期效果：",
            "- 收敛速度提升15%以上",
            "- 解质量改善10%以上",
            "- 算子设计时间减少50%以上",
        ],
    )

    # Slide 15: Summary
    create_content_slide_v2(
        prs,
        "六、总结与后续计划",
        [
            "已完成工作：",
            "- 第三章：生鲜物流配送路径优化模型构建",
            "- 第四章：基于大模型与ALNS算法框架设计",
            " ",
            "后续工作计划：",
            "- 完成实验验证与结果分析",
            "- 完善论文撰写",
            "- 准备毕业答辩",
        ],
    )

    # Slide 16: Thank you
    create_title_slide_v2(prs, "谢谢！", "恳请各位老师批评指正", "")

    prs.save(output_path)
    print("PPT saved to: " + output_path)


if __name__ == "__main__":
    main()
